"""Build all OtoCPA sales PPTX presentations."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Shared Constants ---
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0x9D, 0xC3, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xB0, 0x50)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
DARK_BG2 = RGBColor(0x0F, 0x1A, 0x30)
MID_BLUE = RGBColor(0x18, 0x2C, 0x50)
LAYER_BLUE = RGBColor(0x40, 0x60, 0x90)
HEADER_DARK = RGBColor(0x10, 0x20, 0x40)

W = Inches(13.333)
H = Inches(7.5)
OUT_DIR = os.path.dirname(__file__)


# --- Shared helpers ---
def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def add_bg(slide, color=DARK_BLUE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_logo(slide):
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "OtoCPA"; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = " AI"; r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Segoe UI"


def add_accent_line(slide, y=Inches(1.85)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(2), Pt(4))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()


def add_title(slide, text, y=Inches(1.0), font_size=Pt(34), color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(1.2))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = font_size; p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"


def add_subtitle(slide, text, y=Inches(1.8), color=LIGHT_BLUE, font_size=Pt(22)):
    txBox = slide.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(0.7))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = font_size; p.font.color.rgb = color; p.font.name = "Segoe UI"


def add_body(slide, text, x=Inches(0.7), y=Inches(2.8), w=Inches(11.5), h=Inches(4), font_size=Pt(22), color=WHITE, bold=False, line_spacing=Pt(10)):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = font_size; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.font.bold = bold; p.space_after = line_spacing


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_box(slide, x, y, w, h, fill_color, text_lines, title_size=Pt(18), body_size=Pt(16), title_color=WHITE, body_color=LIGHT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    for i, (text, is_title) in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = title_size if is_title else body_size
        p.font.bold = is_title
        p.font.color.rgb = title_color if is_title else body_color
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)


def title_slide(prs, main_title, subtitle_text, notes_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.5))
    tf = txBox.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "OtoCPA"; r1.font.size = Pt(56); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = " AI"; r2.font.size = Pt(56); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Segoe UI"
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.2), Inches(3), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]; p2.text = main_title; p2.font.size = Pt(28); p2.font.color.rgb = LIGHT_BLUE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.7))
    tf3 = txBox3.text_frame; p3 = tf3.paragraphs[0]; p3.text = subtitle_text; p3.font.size = Pt(20); p3.font.color.rgb = WHITE; p3.font.name = "Segoe UI"; p3.alignment = PP_ALIGN.CENTER
    add_notes(slide, notes_text)
    return slide


def section_slide(prs, number, title, subtitle="", notes=""):
    """Create a section divider slide with a big number."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    # Big number
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(3), Inches(2))
    tf = txBox.text_frame; p = tf.paragraphs[0]; p.text = str(number); p.font.size = Pt(120); p.font.bold = True; p.font.color.rgb = GOLD; p.font.name = "Segoe UI"
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(4.5), Inches(2.2), Inches(8), Inches(1.2))
    tf2 = txBox2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; p2.text = title; p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"
    if subtitle:
        txBox3 = slide.shapes.add_textbox(Inches(4.5), Inches(3.5), Inches(8), Inches(1))
        tf3 = txBox3.text_frame; tf3.word_wrap = True; p3 = tf3.paragraphs[0]; p3.text = subtitle; p3.font.size = Pt(22); p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"
    if notes:
        add_notes(slide, notes)
    return slide


def content_slide(prs, title_text, body_lines, notes="", accent_y=Inches(1.85)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, title_text)
    add_accent_line(slide, y=accent_y)
    if body_lines:
        add_body(slide, "\n".join(body_lines), y=Inches(2.3))
    if notes:
        add_notes(slide, notes)
    return slide


def quote_slide(prs, quote_fr, quote_en="", notes="", title_text=""):
    """A slide with a big quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    if title_text:
        add_title(slide, title_text, font_size=Pt(28), color=GOLD)
        add_accent_line(slide, y=Inches(1.7))
        qy = Inches(2.2)
    else:
        qy = Inches(1.5)
    txBox = slide.shapes.add_textbox(Inches(1.2), qy, Inches(10.5), Inches(3))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f'"{quote_fr}"'; p.font.size = Pt(26); p.font.italic = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    if quote_en:
        p2 = tf.add_paragraph(); p2.text = f'"{quote_en}"'; p2.font.size = Pt(20); p2.font.italic = True; p2.font.color.rgb = LIGHT_BLUE; p2.font.name = "Segoe UI"; p2.space_before = Pt(16)
    if notes:
        add_notes(slide, notes)
    return slide


def two_col_bullets(prs, title_text, left_items, right_items, notes="", left_title="", right_title=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, title_text)
    add_accent_line(slide)

    for col_idx, (items, col_title) in enumerate([(left_items, left_title), (right_items, right_title)]):
        x = Inches(0.7 + col_idx * 6.2)
        y_start = Inches(2.3)
        if col_title:
            txBox = slide.shapes.add_textbox(x, y_start, Inches(5.8), Inches(0.5))
            tf = txBox.text_frame; p = tf.paragraphs[0]; p.text = col_title; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = GOLD; p.font.name = "Segoe UI"
            y_start += Inches(0.55)
        for i, item in enumerate(items):
            y = y_start + Inches(i * 0.55)
            txBox = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.5))
            tf = txBox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(17); r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
            r2 = p.add_run(); r2.text = item; r2.font.size = Pt(17); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    if notes:
        add_notes(slide, notes)
    return slide


# ============================================================
# 1. SALES_SCRIPT.pptx
# ============================================================
def build_sales_script():
    prs = new_prs()

    # Slide 1 - Title
    title_slide(prs,
        "Script de vente --- Mot a mot",
        "Guide complet pour les representants et partenaires",
        "Ce document est un guide interne. Chaque mot est choisi pour maximiser l'impact. "
        "Suivez le script mot a mot lors de vos premiers appels, puis adaptez une fois que vous maitrisez le flux.")

    # Slide 2 - Cold Call Opening
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "1. Ouverture d'appel a froid (30 secondes)")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(2.2), MID_BLUE, [
        ("Script FR:", True),
        ('"Bonjour [prenom], c\'est [votre nom] de OtoCPA. Je ne vous derangerai', False),
        ('pas longtemps. On aide les cabinets CPA du Quebec a eliminer la saisie', False),
        ('manuelle avec l\'IA. Est-ce que la saisie de factures prend beaucoup', False),
        ('de temps dans votre cabinet?"', False),
    ])
    add_body(slide, "Dites leur nom  |  Une seule phrase  |  Une seule question\nNe pitchez pas  |  Si oui: \"Combien d'heures par semaine?\"  |  Ecoutez",
             y=Inches(4.8), font_size=Pt(20), color=GOLD)
    add_notes(slide,
        "Regles: Dites leur nom. Presentez-vous en une phrase. Posez UNE question. "
        "Ne pitchez pas. S'ils disent oui, demandez combien d'heures. Ecoutez. "
        "S'ils disent non ou occupes: proposez un courriel d'une ligne avec un lien de demo de 2 min.")

    # Slide 3 - Voicemail
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "2. Message vocal (20 secondes max)")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(1.8), MID_BLUE, [
        ("Script FR:", True),
        ('"Bonjour [prenom], c\'est [votre nom] de OtoCPA. On aide les cabinets CPA', False),
        ('a eliminer la saisie manuelle --- un de vos collegues economise 160 heures/mois.', False),
        ('Je vous envoie un courriel avec les details. Bonne journee."', False),
    ])
    items = [
        "20 secondes MAXIMUM --- plus long = ils n'ecoutent pas",
        "Nommez un resultat specifique (160 heures, 22 500$, 30 clients)",
        "Dites que vous envoyez un courriel --- puis envoyez-le dans les 5 min",
        "Ne demandez PAS de vous rappeler --- le courriel fait le travail",
    ]
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(4.4 + i * 0.55), Inches(11), Inches(0.5))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(18); r1.font.color.rgb = ORANGE; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "20 secondes maximum. Resultat specifique. Envoyez le courriel dans les 5 minutes.")

    # Slide 4 - Subject Lines
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "3. Objets de courriels qui se font ouvrir")
    add_accent_line(slide)
    subjects = [
        '1.  "160 heures/mois de saisie eliminees --- [NomDuCabinet]"',
        '2.  "Question rapide sur la saisie de donnees chez [NomDuCabinet]"',
        '3.  "TPS/TVQ automatique --- vu par un cabinet CPA du Quebec"',
        '4.  "[Prenom], 2 minutes pour voir comment ca marche?"',
        '5.  "Loi 25 + saisie automatique --- pertinent pour [NomDuCabinet]?"',
        '6.  "De 5 heures a 25 minutes par client"',
    ]
    for i, subj in enumerate(subjects):
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(2.3 + i * 0.65), Inches(11), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        p.text = subj; p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    add_body(slide, "Utilisez leur nom/cabinet  |  Moins de 50 caracteres  |  Utilisez des chiffres\nJamais 'Demo gratuite' ou 'Offre speciale' (filtres anti-spam)",
             y=Inches(6.3), font_size=Pt(18), color=GOLD)
    add_notes(slide, "Personnalisation double le taux d'ouverture. Mobile tronque a 40 caracteres. Utilisez des chiffres.")

    # Slide 5 - Discovery Questions
    two_col_bullets(prs,
        "4. Questions de decouverte (avant de montrer quoi que ce soit)",
        [
            "Combien d'heures de saisie par semaine?",
            "Comment recevez-vous les documents?",
            "Plus gros probleme en saison des impots?",
            "Deja eu un doublon qui est passe?",
            "Outil actuel? Papier? Excel? Logiciel?",
            "QuickBooks Desktop ou Online?",
        ],
        [
            "Combien de clients actifs?",
            "Combien de personnes font de la saisie?",
            "Que feriez-vous avec 20h de plus/semaine?",
            "Loi 25 vous preoccupe?",
            "Missions de verification ou d'examen?",
            "Erreur de saisie vs Revenu Quebec?",
        ],
        notes="Posez. Taisez-vous. Prenez des notes. Ne les interrompez jamais. "
              "Ecrivez leur douleur mot pour mot --- vous la reutiliserez sur la diapo 2.",
        left_title="Douleur", right_title="Contexte")

    # Slide 6 - Transition to Demo
    quote_slide(prs,
        "Merci, c'est tres utile. Vous avez mentionne [repeter leur probleme dans LEURS mots]. "
        "C'est exactement ce que OtoCPA regle. Est-ce que je peux vous montrer en 15 minutes?",
        notes="Repetez LEURS mots, pas les votres. Demandez la permission de montrer la demo. "
              "Fixez l'attente de temps: 15 minutes.",
        title_text="5. Transition vers la demo")

    # Slide 7 - Demo Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "6. Narration de demo (resume)")
    add_accent_line(slide)
    demo_steps = [
        ("1", "Document traite\nautomatiquement", "Photo vers compte GL", ACCENT_BLUE),
        ("2", "Detection de\nfraude en direct", "13 regles en temps reel", ORANGE),
        ("3", "Resume de\nproduction FPZ-500", "Pre-remplissage RQ", GREEN),
        ("4", "Module d'audit\nCAS 315-700", "Dossiers de travail", GOLD),
    ]
    for i, (num, title, desc, color) in enumerate(demo_steps):
        x = Inches(0.5 + i * 3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.95), Inches(3.5))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(2.7), Inches(0.8), Inches(0.8))
        circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
        tf = circle.text_frame; p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE if color != GOLD else DARK_BLUE; p.font.name = "Segoe UI"
        # Title
        txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(3.7), Inches(2.65), Inches(1))
        tf2 = txBox.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; p2.text = title; p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(20); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"
        # Desc
        txBox3 = slide.shapes.add_textbox(x + Inches(0.15), Inches(4.9), Inches(2.65), Inches(0.6))
        tf3 = txBox3.text_frame; tf3.word_wrap = True; p3 = tf3.paragraphs[0]; p3.text = desc; p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Pt(16); p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"
    add_notes(slide, "Dites: 'Je vais vous montrer 4 choses en 15 minutes.' Puis listez-les. 'On commence.'")

    # Slide 8 - Price Conversation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "7. La conversation sur le prix")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2), ACCENT_BLUE, [
        ("Quand ils demandent le prix:", True),
        ('"Le forfait Professionnel, qui convient a la', False),
        ("plupart des cabinets de votre taille, est a", False),
        ('249$/mois. Plus 500-1000$ d\'installation."', False),
    ])
    add_box(slide, Inches(6.6), Inches(2.3), Inches(6.1), Inches(2), MID_BLUE, [
        ("Puis: ARRETEZ DE PARLER", True),
        ("Laissez le silence travailler.", False),
        ("Le prochain a parler perd la negociation.", False),
    ], title_color=GOLD)
    rules = [
        'Ne dites JAMAIS "seulement" ou "juste" --- ca sonne apologetique',
        "Ancrez sur le ROI, pas le prix: le prix est 249$, la valeur est 24 000$",
        "S'ils disent 'c'est cher': faites le calcul avec LEURS chiffres",
    ]
    for i, rule in enumerate(rules):
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(4.8 + i * 0.6), Inches(11.5), Inches(0.55))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(18); r1.font.color.rgb = ORANGE; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = rule; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "Dites le prix avec confiance. Arretez de parler. Ancrez sur le ROI. Silence = outil de vente.")

    # Slide 9 - How to Close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "8. Techniques de cloture")
    add_accent_line(slide)
    closes = [
        ("Cloture directe", '"On commence avec un pilote de 30 jours?\nJe peux configurer 5 clients cette semaine."', GREEN, "Pour ceux qui semblent prets"),
        ("Cloture calendrier", '"Mardi prochain a 10h, ca marcherait\npour une demo de 15 minutes?"', ACCENT_BLUE, "Pour ceux qui ont besoin de temps"),
        ("Cloture courriel", '"Je vous envoie un resume. Est-ce que\nje peux vous recontacter jeudi?"', GOLD, "Pour ceux qui sont occupes"),
    ]
    for i, (name, script, color, desc) in enumerate(closes):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(3.9), Inches(3.8))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(0.5))
        tf = txBox.text_frame; p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        txBox2 = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(1.8))
        tf2 = txBox2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; p2.text = script; p2.font.size = Pt(17); p2.font.italic = True; p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER
        txBox3 = slide.shapes.add_textbox(x + Inches(0.2), Inches(5.2), Inches(3.5), Inches(0.5))
        tf3 = txBox3.text_frame; p3 = tf3.paragraphs[0]; p3.text = desc; p3.font.size = Pt(15); p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"; p3.alignment = PP_ALIGN.CENTER
    add_body(slide, "TOUJOURS proposer une prochaine etape concrete  |  TOUJOURS obtenir une date",
             y=Inches(6.5), font_size=Pt(20), color=GOLD)
    add_notes(slide, "Toujours proposer une prochaine etape concrete. Toujours obtenir une date. "
              "'La semaine prochaine' n'est pas une date. 'Mardi a 14h' est une date.")

    # Slide 10 - Follow-Up Sequence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "9. Sequence de relance")
    add_accent_line(slide)
    followups = [
        ("Jour 1", "Resume + prochaines etapes\n(dans les 2 heures)", ACCENT_BLUE, "OBLIGATOIRE"),
        ("Jour 3", "Courriel court:\n'La demo de 15 min, ca vous dit?'", LAYER_BLUE, "Si pas de reponse"),
        ("Jour 7", "Derniere question:\n'Projet 2026 ou plus tard?'", LAYER_BLUE, "Si pas de reponse"),
        ("Jour 14", "Article de valeur\nAUCUNE demande", GREEN, "Valeur pure"),
    ]
    for i, (day, desc, color, label) in enumerate(followups):
        x = Inches(0.5 + i * 3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.95), Inches(3.5))
        box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = day; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = GOLD; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(17); p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
        p3 = tf.add_paragraph(); p3.text = label; p3.font.size = Pt(14); p3.font.bold = True; p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(12)
    add_body(slide, "Jamais plus de 4 relances  |  Chaque relance plus courte que la precedente\nLe jour 14 n'a AUCUNE demande --- c'est de la valeur pure",
             y=Inches(6.2), font_size=Pt(18), color=GOLD)
    add_notes(slide, "Jour 1 est obligatoire (2h). 4 relances max. Chaque relance plus courte. "
              "Jour 14 = valeur pure, aucune demande. C'est ce qui les fait se souvenir de vous 3 mois plus tard.")

    # Slide 11 - Qualifying Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "10. Liste de qualification")
    add_accent_line(slide)
    checks = [
        ("Cabinet CPA au Quebec?", "Obligatoire"),
        ("Au moins 5 clients actifs?", "En dessous de 5, le ROI est faible"),
        ("QuickBooks (Desktop ou Online)?", "Pas obligatoire mais renforce la demo"),
        ("Decideur present dans la conversation?", "Sinon, obtenez un RDV avec lui"),
        ("Douleur autour de la saisie ou conformite?", "Pas de douleur = pas de vente"),
        ("Budget pour 500-1000$ d'installation?", "Sinon, offrez le pilote d'abord"),
    ]
    for i, (check, detail) in enumerate(checks):
        y = Inches(2.3 + i * 0.75)
        txBox = slide.shapes.add_textbox(Inches(0.9), y, Inches(11), Inches(0.65))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(18); r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = check + "  "; r2.font.size = Pt(19); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
        r3 = p.add_run(); r3.text = detail; r3.font.size = Pt(16); r3.font.color.rgb = LIGHT_BLUE; r3.font.name = "Segoe UI"
    add_body(slide, "4+ cases cochees = demo complete  |  Moins = focus sur la relation + one-pager",
             y=Inches(7.0), font_size=Pt(20), color=GOLD)
    add_notes(slide, "Si 4+ boxes cochees, procedez avec la demo complete. "
              "Si moins, concentrez-vous sur la relation et envoyez le one-pager.")

    out = os.path.join(OUT_DIR, "SALES_SCRIPT.pptx")
    prs.save(out)
    print(f"  SALES_SCRIPT.pptx --- {len(prs.slides)} slides")


# ============================================================
# 2. OBJECTIONS.pptx
# ============================================================
def build_objections():
    prs = new_prs()

    # Title
    title_slide(prs,
        "Guide de traitement des objections",
        "Reconnaitre. Recadrer. Repondre. Avancer.",
        "Regle #1: Chaque objection est une question deguisee. Trouvez la question. Repondez-y.\n"
        "Regle #2: Ne discutez jamais. Reconnaissez, recadrez, repondez.\n"
        "Regle #3: Apres votre reponse, posez une question. Avancez.")

    # Master Framework slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Methode universelle: Reconnaitre-Recadrer-Repondre")
    add_accent_line(slide)
    steps = [
        ("1. Reconnaitre", '"Je comprends."\n"C\'est une bonne question."', ACCENT_BLUE),
        ("2. Recadrer", '"Ce que vous demandez\nvraiment, c\'est..."', LAYER_BLUE),
        ("3. Repondre", "Faits, chiffres,\nou histoire.", GREEN),
        ("4. Avancer", '"Est-ce que ca repond?\nSi oui, [prochaine etape]"', GOLD),
    ]
    for i, (name, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(2.95), Inches(2.5))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(17); p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
    add_notes(slide, "C'est la methode pour TOUTE objection non listee. Reconnaitre, recadrer, repondre, avancer.")

    objections = [
        ("On utilise deja TaxDome / Karbon / CaseWare",
         "Ils ont deja investi du temps et de l'argent ailleurs.",
         '"Comment ca se passe avec [leur outil]? Est-ce que la saisie est vraiment automatisee?"',
         "OtoCPA ne remplace pas [outil]. Il s'ajoute en amont.\n[Outil] gere le flux de travail. OtoCPA elimine la saisie.\nLes deux peuvent coexister.",
         '"Ca vaudrait la peine de voir comment les deux fonctionnent ensemble?"'),
        ("C'est trop cher",
         "Ils ne voient pas assez de valeur. Ou doivent justifier a un associe.",
         '"Je comprends. Regardons les chiffres ensemble."',
         "[X] clients x [Y]h de saisie = [Z]h/mois\nA 150$/h = [$$] de capacite\nOtoCPA = 249$/mois = 1.7% du temps recupere\n= Rentable en moins de 2 heures",
         '"On offre un pilote de 30 jours gratuit. Vous testez, vous mesurez, vous decidez."'),
        ("On n'a pas le temps d'implanter",
         "Ils sont deja submerges. Un logiciel de plus les effraie.",
         '"C\'est justement pour ca que OtoCPA existe."',
         "Installation: 30 minutes a distance. VOUS faites tout.\nInterface simple: file d'attente, un clic pour approuver.\nAucune formation complexe requise.",
         '"30 minutes cette semaine, ca serait faisable?"'),
        ("Nos clients ne l'utiliseront pas",
         "Peur du changement cote client.",
         '"Vos clients n\'ont rien a apprendre."',
         "Seul changement: photo au lieu de scanner/enveloppe.\nCode QR par client --- pas d'app, pas de mot de passe.\nCourriels et PDF traites de la meme facon.",
         '"Avez-vous des clients technos pour le pilote?"'),
        ("Et si vous fermez?",
         "Mefiance envers les startups.",
         '"C\'est l\'avantage principal de OtoCPA."',
         "Vos donnees = sur VOTRE serveur. SQLite lisible partout.\nSi OtoCPA disparait, vos donnees restent.\nComparez: TaxDome/Karbon = serveurs Amazon aux USA.",
         '"La question des donnees est votre principale preoccupation?"'),
        ("On va y penser",
         "Pas prets a decider. Veulent eviter de dire non en face.",
         '"Bien sur. Est-ce que je peux poser une question?"',
         "\"Qu'est-ce qui vous ferait pencher d'un cote ou de l'autre?\"\nSi flou: \"Le prix, le timing, ou la confiance dans la techno?\"\nIls choisiront. Maintenant vous savez quoi adresser.",
         '"Est-ce que je peux vous recontacter [jour precis]?"'),
        ("Peut-on avoir un rabais?",
         "SIGNAL D'ACHAT. Ils veulent acheter mais veulent un deal.",
         'NE JAMAIS baisser le prix mensuel.',
         "Option A: +5 clients configures dans les frais d'installation\nOption B: Engagement 12 mois = installation offerte\nOption C: Premier mois a 50% comme client fondateur",
         '"Le prix reflete la valeur. 249$ pour recuperer 15 000$."'),
        ("L'IA fait des erreurs?",
         "Peur que l'IA ruine les livres de leurs clients.",
         '"C\'est la bonne question."',
         "Taxes = JAMAIS par IA. Math pur, librairie Decimal, 0 erreur.\nIA = lecture de documents + suggestions GL seulement.\nChaque transaction approuvee par un humain.\nFile de revision pour confiance basse.",
         '"L\'IA peut suggerer une erreur. Mais ne peut JAMAIS la publier."'),
    ]

    for obj_text, real_meaning, response_start, body, close_q in objections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_logo(slide)
        # Objection as title
        add_title(slide, f'"{obj_text}"', font_size=Pt(28), color=ORANGE)
        add_accent_line(slide, y=Inches(1.75))
        # Real meaning
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(12), Inches(0.5))
        tf = txBox.text_frame; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "Ils pensent vraiment: "; r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = real_meaning; r2.font.size = Pt(17); r2.font.italic = True; r2.font.color.rgb = LIGHT_BLUE; r2.font.name = "Segoe UI"
        # Response box
        add_box(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(1), ACCENT_BLUE, [
            ("Ouvrez avec:", True), (response_start, False)
        ], title_size=Pt(15), body_size=Pt(16))
        # Body/arguments
        add_body(slide, body, x=Inches(0.7), y=Inches(4.1), font_size=Pt(18), color=WHITE)
        # Closing question
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(12), Inches(0.5))
        tf2 = txBox2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]
        r1 = p2.add_run(); r1.text = "Concluez: "; r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
        r2 = p2.add_run(); r2.text = close_q; r2.font.size = Pt(17); r2.font.italic = True; r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
        add_notes(slide, f'Objection: "{obj_text}"\nIls pensent: {real_meaning}\n'
                  f'Reponse: {response_start}\n{body}\nConcluez: {close_q}')

    # CPA Approval + Too Small/Big slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "Autres objections courantes")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2.2), MID_BLUE, [
        ('"Approuve par l\'Ordre des CPA?"', True),
        ("L'Ordre n'approuve AUCUN logiciel.", False),
        ("Ni QBO, ni CaseWare, ni aucun autre.", False),
        ("Exigence: jugement pro + piste d'audit.", False),
        ("OtoCPA = revue humaine + audit complet.", False),
    ], title_color=ORANGE)
    add_box(slide, Inches(6.6), Inches(2.3), Inches(6.1), Inches(1), MID_BLUE, [
        ('"On est trop petit" (<10 clients)', True),
        ("Essentiel: 99$/mois, 10 clients, 3 users. Rentable en 1h.", False),
    ], title_color=ORANGE)
    add_box(slide, Inches(6.6), Inches(3.6), Inches(6.1), Inches(1), MID_BLUE, [
        ('"On est trop gros" (>75 clients)', True),
        ("Entreprise: 999$/mois, illimite + audit CAS complet.", False),
    ], title_color=ORANGE)
    add_box(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(1.2), MID_BLUE, [
        ('"On vient de renouveler avec [concurrent]"', True),
        ("OtoCPA s'ajoute en amont, ne remplace pas. Notez la date de fin.", False),
        ("Offrez le pilote gratuit maintenant --- ca ne touche pas leur contrat.", False),
    ], title_color=ORANGE)
    add_notes(slide, "L'Ordre n'approuve aucun logiciel. Trop petit = Essentiel 99$. "
              "Trop gros = Entreprise 999$. Concurrent renouvele = pilote gratuit maintenant.")

    out = os.path.join(OUT_DIR, "OBJECTIONS.pptx")
    prs.save(out)
    print(f"  OBJECTIONS.pptx --- {len(prs.slides)} slides")


# ============================================================
# 3. ONE_PAGER_FR.pptx
# ============================================================
def build_one_pager():
    prs = new_prs()

    # Title
    title_slide(prs,
        "La comptabilite intelligente pour les cabinets CPA du Quebec",
        "Fiche-resume | One-Pager",
        "OtoCPA elimine la saisie manuelle de factures. Vos clients envoient une photo. "
        "L'IA extrait les donnees. Votre equipe approuve en un clic. Les taxes sont calculees automatiquement.")

    # 3 Key Advantages
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "3 avantages cles")
    add_accent_line(slide)
    advantages = [
        ("160 h/mois", "recuperees", "De 5h de saisie/client a 25 min.\n30 clients = 24 000$/mois a 150$/h.\nROI: 9 600%.", ACCENT_BLUE),
        ("Zero erreur", "de taxe", "TPS/TVQ/TVH par regles math, pas IA.\nCalcul parallele (jamais en cascade).\nFPZ-500 pre-rempli. 98 tests: 0 echec.", GREEN),
        ("Donnees", "chez vous", "Serveur local. SQLite. Loi 25 conforme.\nJamais dans le cloud.\nAucun serveur americain.", GOLD),
    ]
    for i, (big, small, desc, color) in enumerate(advantages):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(3.9), Inches(4.5))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = big; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = small; p2.font.size = Pt(24); p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = ""; p3.font.size = Pt(8)
        for line in desc.split("\n"):
            p4 = tf.add_paragraph(); p4.text = line; p4.font.size = Pt(16); p4.font.color.rgb = LIGHT_BLUE; p4.font.name = "Segoe UI"; p4.alignment = PP_ALIGN.CENTER; p4.space_after = Pt(4)
    add_notes(slide, "Trois avantages cles: 160h/mois recuperees, zero erreur de taxe, donnees locales (Loi 25).")

    # Features Grid
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Fonctionnalites incluses")
    add_accent_line(slide)
    features = [
        ("Extraction IA", "Photos, PDF, courriels"),
        ("Detection de fraude", "13 regles deterministes"),
        ("Rapprochement bancaire", "Desjardins, BMO, TD, RBC, BN"),
        ("Module d'audit", "CAS 315-700, CSQC 1"),
        ("Calendrier de production", "Echeances TPS/TVQ + rappels"),
        ("Suivi du temps", "Entrees + facturation TPS/TVQ"),
        ("Communications clients", "Messages bilingues par IA"),
        ("Microsoft 365", "Courriel, SharePoint, OneDrive"),
        ("Bilingue FR/EN", "100% traduit, terminologie CPA"),
    ]
    for i, (name, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(2.3 + row * 1.55)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(1.3))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(19); p.font.bold = True; p.font.color.rgb = GOLD; p.font.name = "Segoe UI"
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(16); p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"
    add_notes(slide, "9 fonctionnalites principales couvrant extraction, fraude, rapprochement, audit, calendrier, temps, comms, M365, bilingue.")

    # Pricing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "Forfaits")
    add_accent_line(slide)
    plans = [
        ("Essentiel", "99$", "10 clients | 3 users", "CPA solo", ACCENT_BLUE),
        ("Professionnel", "249$", "30 clients | 5 users", "Petit cabinet", RGBColor(0x2E, 0x75, 0xB6)),
        ("Cabinet", "499$", "75 clients | 15 users", "Cabinet en croissance", LAYER_BLUE),
        ("Entreprise", "999$", "Illimite", "Grand cabinet + audit", GOLD),
    ]
    for i, (name, price, capacity, target, color) in enumerate(plans):
        x = Inches(0.5 + i * 3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.95), Inches(3.5))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = price; p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = "/mois"; p3.font.size = Pt(16); p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"; p3.alignment = PP_ALIGN.CENTER
        p4 = tf.add_paragraph(); p4.text = capacity; p4.font.size = Pt(15); p4.font.color.rgb = WHITE; p4.font.name = "Segoe UI"; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(12)
        p5 = tf.add_paragraph(); p5.text = target; p5.font.size = Pt(14); p5.font.italic = True; p5.font.color.rgb = LIGHT_BLUE; p5.font.name = "Segoe UI"; p5.alignment = PP_ALIGN.CENTER
    add_body(slide, "Installation: 500$ - 1 000$ (une fois)  |  Pilote gratuit de 30 jours",
             y=Inches(6.2), font_size=Pt(20), color=GOLD)
    add_notes(slide, "4 forfaits: Essentiel 99$, Professionnel 249$, Cabinet 499$, Entreprise 999$. Installation unique 500-1000$. Pilote 30 jours gratuit.")

    # Security
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Securite et Loi 25")
    add_accent_line(slide)
    security_items = [
        "Donnees hebergees sur votre serveur local --- jamais dans le cloud",
        "Authentification bcrypt avec sessions securisees",
        "Piste d'audit complete (qui, quoi, quand)",
        "Sauvegardes automatiques quotidiennes",
        "Acces distant securise via Cloudflare Tunnel (HTTPS)",
        "Aucun acces de OtoCPA a vos donnees clients",
    ]
    for i, item in enumerate(security_items):
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.3 + i * 0.7), Inches(10), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(20); r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(22); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "Securite: serveur local, bcrypt, audit complet, sauvegardes auto, Cloudflare, zero acces OtoCPA.")

    # Contact / CTA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1))
    tf = txBox.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "OtoCPA"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = " AI"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Segoe UI"
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(2.8), Inches(3), Pt(3)); s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
    add_body(slide, "support@otocpa.com\notocpa.com\n\nReservez une demo de 15 minutes\nOn vous montre le logiciel avec vos vrais documents.",
             y=Inches(3.2), font_size=Pt(24), color=WHITE)
    add_notes(slide, "Coordonnees: support@otocpa.com, otocpa.com. Demo de 15 minutes avec vrais documents.")

    out = os.path.join(OUT_DIR, "ONE_PAGER_FR.pptx")
    prs.save(out)
    print(f"  ONE_PAGER_FR.pptx --- {len(prs.slides)} slides")


# ============================================================
# 4. EMAIL_TEMPLATES.pptx
# ============================================================
def build_email_templates():
    prs = new_prs()

    # Title
    title_slide(prs,
        "Modeles de courriels",
        "Copiez. Personnalisez. Envoyez.",
        "Ne jamais envoyer sans personnaliser les [crochets]. "
        "Chaque courriel doit sembler ecrit specifiquement pour le destinataire.")

    # Rules slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Regles generales pour TOUS les courriels")
    add_accent_line(slide)
    rules = [
        ("< 150 mots", "pour le cold outreach, < 200 pour les relances"),
        ("1 seul CTA", "par courriel --- pas de video + demo + site web"),
        ("Personnalisez", "tout ce qui est entre [crochets]"),
        ("Mardi-Jeudi", "entre 8h et 10h --- meilleur taux d'ouverture"),
        ("Prenom", 'jamais "Dear" ou "A qui de droit"'),
        ("Pas de PJ", "en cold email --- declenchent les filtres anti-spam"),
        ("Re:", "dans l'objet des relances pour la continuite"),
        ("Telephone", "inclus --- certains preferent appeler"),
    ]
    for i, (bold_part, rest) in enumerate(rules):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(2.3 + row * 1.1)
        txBox = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.9))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = bold_part + "  "; r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "8 regles universelles: court, 1 CTA, personnalise, mardi-jeudi 8-10h, prenom, pas de PJ, Re:, telephone.")

    # Template slides
    templates = [
        ("1. Premier contact (Cold Outreach)",
         "Objet: '[Prenom], 160 heures/mois de saisie eliminees?'",
         [
            "Bonjour [Prenom],",
            "",
            "Je m'appelle [votre nom], je travaille chez OtoCPA.",
            "",
            "On aide les cabinets CPA du Quebec a eliminer la saisie",
            "manuelle de factures avec l'IA. Le calcul TPS/TVQ est",
            "100% base sur des regles --- zero erreur.",
            "",
            "Pour [X] clients comme [NomDuCabinet], ca represente",
            "environ [Y] heures de saisie eliminees par mois.",
            "",
            "Demo de 15 minutes cette semaine?",
         ],
         "Premier email. Court, personnalise, un seul CTA: la demo de 15 min."),

        ("2. Relance Jour 3 (pas de reponse)",
         "Objet: 'Re: [Prenom], 160 heures/mois de saisie eliminees?'",
         [
            "Bonjour [Prenom],",
            "",
            "Je voulais m'assurer que mon courriel s'est bien rendu.",
            "",
            "En un mot: OtoCPA automatise la saisie comptable.",
            "Photo de facture -> extraction auto -> clic pour approuver.",
            "",
            "15 minutes pour voir --- ca vous dit?",
         ],
         "Plus court que le premier. Utilisez Re: dans l'objet."),

        ("3. Relance Jour 7 (derniere question)",
         "Objet: 'Derniere question, [Prenom]'",
         [
            "Bonjour [Prenom],",
            "",
            "Je ne veux pas encombrer votre boite de reception.",
            "",
            "Est-ce que l'automatisation de la saisie est un projet",
            "pour 2026 chez [NomDuCabinet], ou pour plus tard?",
            "",
            "Si plus tard, je me note de vous recontacter.",
         ],
         "Ton respectueux. Offrez de recontacter plus tard. Laissez une porte ouverte."),

        ("4. Suivi post-demo",
         "Objet: 'Merci pour la demo, [Prenom] --- resume + prochaines etapes'",
         [
            "Bonjour [Prenom], merci pour votre temps.",
            "",
            "CE QU'ON A VU:",
            "- Extraction auto (photo -> GL en 10s)",
            "- Detection de fraude (13 regles)",
            "- Pre-remplissage FPZ-500",
            "",
            "VOS CHIFFRES:",
            "- [X] clients, [Y]h de saisie/mois",
            "- Valeur: [Z]$/mois a 150$/h",
            "- Forfait recommande: [plan] a [prix]$/mois",
            "",
            "PROCHAINE ETAPE: [pilote/demo equipe/rappel]",
         ],
         "Envoyez dans les 2 heures. Incluez LEURS chiffres specifiques."),

        ("5. Proposition avec prix",
         "Objet: 'Proposition OtoCPA pour [NomDuCabinet]'",
         [
            "Suite a notre conversation, voici notre proposition.",
            "",
            "FORFAIT: [Professionnel/Cabinet/Entreprise]",
            "- [249/499/999]$/mois",
            "- [30/75/illimite] clients",
            "- Installation: [500-1000]$ (une fois)",
            "",
            "INCLUS: Extraction IA, TPS/TVQ auto, fraude,",
            "rapprochement bancaire, FPZ-500, bilingue, Loi 25",
            "",
            "PILOTE GRATUIT: 30 jours, 5 clients, 0 obligation",
         ],
         "Email formel de proposition. Resumez tout clairement. Un seul CTA: commencer le pilote."),

        ("6. Renouvellement (30 jours avant)",
         "Objet: 'Renouvellement OtoCPA --- expire le [date]'",
         [
            "Votre licence [forfait] expire le [date].",
            "",
            "POUR RENOUVELER:",
            "1. Nouvelle cle de licence pour [12/24] mois",
            "2. Activer dans le tableau de bord",
            "3. C'est tout.",
            "",
            "VOTRE FORFAIT ACTUEL:",
            "- [Plan] a [prix]$/mois",
            "- [X]/[Y] clients actifs",
            "",
            "Changement de forfait? C'est le bon moment.",
         ],
         "Envoyez 30 jours avant l'expiration. Grace de 30 jours mais fonctionnalites limitees."),

        ("7. Win-Back (prospect qui a dit 'pas maintenant')",
         "Objet: '[Prenom], les choses ont change depuis [mois]'",
         [
            "On s'est parle en [mois]. Le timing n'etait pas bon.",
            "",
            "Depuis, on a ajoute:",
            "- [Fonctionnalite 1 pertinente pour eux]",
            "- [Fonctionnalite 2]",
            "- [Fonctionnalite 3]",
            "",
            "Et on a maintenant [X] cabinets CPA au Quebec.",
            "",
            "Le timing est meilleur? 15 min pour les nouveautes?",
            "Sinon, je me note de recontacter dans [3/6] mois.",
         ],
         "Ton decontracte. Montrez du progres. Laissez une porte ouverte."),
    ]

    for title, subject, lines, notes in templates:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_logo(slide)
        add_title(slide, title, font_size=Pt(28))
        add_accent_line(slide, y=Inches(1.7))
        # Subject line
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(0.5))
        tf = txBox.text_frame; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = subject; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        # Email body in a box
        email_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.6), Inches(11.9), Inches(4.5))
        email_box.fill.solid(); email_box.fill.fore_color.rgb = MID_BLUE; email_box.line.color.rgb = ACCENT_BLUE; email_box.line.width = Pt(1)
        tf2 = email_box.text_frame; tf2.word_wrap = True
        for i, line in enumerate(lines):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = line; p2.font.size = Pt(15); p2.font.name = "Segoe UI"
            p2.font.color.rgb = WHITE if not line.startswith(("[", "CE ", "VOS", "PRO", "FOR", "INC", "PIL", "POU")) else GOLD
            if line.startswith(("[", "CE ", "VOS", "PRO", "FOR", "INC", "PIL", "POU")):
                p2.font.bold = True
        add_notes(slide, notes)

    out = os.path.join(OUT_DIR, "EMAIL_TEMPLATES.pptx")
    prs.save(out)
    print(f"  EMAIL_TEMPLATES.pptx --- {len(prs.slides)} slides")


# ============================================================
# 5. DEMO_SCRIPT.pptx
# ============================================================
def build_demo_script():
    prs = new_prs()

    # Title
    title_slide(prs,
        "Script de demo en direct --- 15 minutes",
        "Guide pas a pas pour le presentateur",
        "Public: Proprietaire/associe de cabinet CPA + potentiellement leur equipe.\n"
        "Pre-requis: OtoCPA installe avec demo data. Navigateur ouvert sur 127.0.0.1:8787/login.")

    # Pre-Demo Setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "Preparation pre-demo (10 min avant)")
    add_accent_line(slide)
    setup_steps = [
        ("1. Demo mode", 'otocpa.config.json: "demo_mode": true'),
        ("2. Charger donnees", "python scripts/load_demo_data.py (50 docs, 5 clients)"),
        ("3. Redemarrer", "sc stop/start OtoCPA ou python review_dashboard.py"),
        ("4. Navigateur", "127.0.0.1:8787/login, connecte, FR, onglet filing summary pret"),
        ("5. Nettoyer", "Fermer email, Slack, notifications. Bureau propre."),
    ]
    for i, (step, desc) in enumerate(setup_steps):
        y = Inches(2.3 + i * 0.9)
        txBox = slide.shapes.add_textbox(Inches(0.9), y, Inches(11.5), Inches(0.8))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = step + "  "; r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_body(slide, "5 clients demo: MARCEL | BOLDUC | DENTAIRE | BOUTIQUE | TECHLAVAL\nChacun a 10 documents incluant: factures, fraude, confiance basse, repas, banque",
             y=Inches(6.8), font_size=Pt(16), color=LIGHT_BLUE)
    add_notes(slide, "Faites cette preparation 10 minutes avant l'appel. "
              "5 clients demo pre-configures avec 50 documents couvrant tous les scenarios.")

    # Timeline overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Chronologie de la demo --- 15 minutes")
    add_accent_line(slide)
    timeline = [
        ("0:00-2:00", "Probleme + Setup", "Question sur les heures de saisie.\nLes 4 choses a montrer.", ACCENT_BLUE),
        ("2:00-5:00", "Extraction auto", "Document photo -> GL en 10s.\nConfiance, TPS/TVQ, Approuver.", GREEN),
        ("5:00-8:00", "Detection de fraude", "Document signale. 13 regles.\nOverride = role manager + raison.", ORANGE),
        ("8:00-11:00", "FPZ-500 / RQ", "Resume de production.\nLignes 103-209 pre-remplies.", ACCENT_BLUE),
        ("11:00-14:00", "Module d'audit", "CAS 315-700. Dossiers de travail.\nEtats financiers. (Si pertinent)", LAYER_BLUE),
        ("14:00-15:00", "Cloture", "ROI personnalise.\nPilote 30 jours gratuit.", GOLD),
    ]
    for i, (time, title, desc, color) in enumerate(timeline):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(2.3 + row * 2.4)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.1))
        box.fill.solid(); box.fill.fore_color.rgb = MID_BLUE; box.line.color.rgb = color; box.line.width = Pt(2)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = time; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Segoe UI"
        p2 = tf.add_paragraph(); p2.text = title; p2.font.size = Pt(20); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.font.name = "Segoe UI"
        for line in desc.split("\n"):
            p3 = tf.add_paragraph(); p3.text = line; p3.font.size = Pt(14); p3.font.color.rgb = LIGHT_BLUE; p3.font.name = "Segoe UI"
    add_notes(slide, "6 segments de la demo. Suivez le chronometrage strictement. 15 minutes = credibilite.")

    # Minute 0-2: Problem + Setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "0:00-2:00 --- Probleme et cadrage")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(2), MID_BLUE, [
        ("Ouvrez avec la question:", True),
        ('"Combien d\'heures par semaine votre equipe passe-t-elle', False),
        ('a saisir des factures dans QuickBooks?"', False),
        ("", False),
        ("Attendez. Notez. Puis faites le calcul a voix haute.", True),
    ])
    add_body(slide, 'Annoncez les 4 choses a montrer:\n1. Document traite automatiquement\n2. Detection de fraude en direct\n3. Resume de production FPZ-500\n4. Module d\'audit (si pertinent)\n\n"On commence."',
             y=Inches(4.6), font_size=Pt(20), color=WHITE)
    add_notes(slide, "Posez la question. Attendez. Notez la reponse. Faites le math a voix haute: "
              "[X]h x 150$/h = [$$]/semaine. Annoncez les 4 segments. 'On commence.'")

    # Minute 2-5: Auto Extraction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "2:00-5:00 --- Extraction automatique")
    add_accent_line(slide)
    actions = [
        ("MONTRER", "File d'attente --- chaque ligne = un document lu et classe"),
        ("CLIQUER", "Document 'Ready' pour client MARCEL"),
        ("POINTER", "Fournisseur | Montant | Date | Compte GL | TPS 5% | TVQ 9,975%"),
        ("INSISTER", "TPS/TVQ en parallele, PAS en cascade --- correct au Quebec"),
        ("SCROLLER", "Score de confiance: >85% = recommande | <85% = revision"),
        ("CLIQUER", '"Approuver" --- un clic, c\'est fait, pret pour QBO'),
    ]
    for i, (action, desc) in enumerate(actions):
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.3 + i * 0.7), Inches(12), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{action}  "; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_body(slide, '"De la photo a l\'approbation: 10 secondes. Comparez a 5 minutes de saisie manuelle."',
             y=Inches(6.6), font_size=Pt(20), color=GOLD, bold=True)
    add_notes(slide, "Montrez la file d'attente. Cliquez sur MARCEL. Pointez chaque champ extrait. "
              "Insistez sur le calcul parallele TPS/TVQ. Montrez le score de confiance. Cliquez Approuver. "
              "Point cle: 10 secondes vs 5 minutes.")

    # Minute 5-8: Fraud Detection
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "5:00-8:00 --- Detection de fraude")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(1), ACCENT_BLUE, [
        ('"Maintenant, je vais vous montrer quelque chose qu\'aucun autre logiciel ne fait."', True),
    ], title_size=Pt(22))
    fraud_actions = [
        ("FILTRER", "par status 'NeedsReview' --- trouver document BOLDUC"),
        ("CLIQUER", "sur le document signale. Scroller aux alertes."),
        ("LIRE", "les flags: Doublon exact | Transaction week-end | Anomalie montant"),
        ("MONTRER", "Override = role manager + raison ecrite = piste d'audit"),
    ]
    for i, (action, desc) in enumerate(fraud_actions):
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.6 + i * 0.7), Inches(12), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{action}  "; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = ORANGE; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_body(slide, '"13 regles. Toutes basees sur des regles, pas IA.\nSeverite CRITIQUE ou HAUTE = auto-approbation impossible. Point final."',
             y=Inches(6.3), font_size=Pt(20), color=GOLD, bold=True)
    add_notes(slide, "Dites: 'Quelque chose qu'aucun autre logiciel ne fait.' Filtrez NeedsReview. "
              "Montrez les flags. Override = manager + raison. 13 regles, pas IA.")

    # Minute 8-11: Filing Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "8:00-11:00 --- Resume de production et Revenu Quebec")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(1), GREEN, [
        ("Naviguer a /filing_summary --- le moment le plus impressionnant", True),
    ], title_size=Pt(20))
    lines_rq = [
        ("Ligne 103", "Ventes et fournitures taxables"),
        ("Ligne 106", "TPS percue"),
        ("Ligne 108", "Total TPS percue"),
        ("Ligne 205", "TVQ percue"),
        ("Ligne 207", "CTI reclames"),
        ("Ligne 209", "RTI reclames"),
    ]
    for i, (line, desc) in enumerate(lines_rq):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(3.6 + row * 0.7)
        txBox = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.6))
        tf = txBox.text_frame; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{line}  "; r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(18); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_body(slide, "Inclut: repas 50% | assurance 9% non recuperable | exemptions\nMontrez le PDF telechargeable | Si methode rapide: montrez le toggle par client",
             y=Inches(5.8), font_size=Pt(18), color=LIGHT_BLUE)
    add_body(slide, '"Ce resume, a la main, ca prend combien de temps? Avec OtoCPA, c\'est instantane."',
             y=Inches(6.6), font_size=Pt(20), color=GOLD, bold=True)
    add_notes(slide, "Naviguez a /filing_summary. Montrez les lignes FPZ-500 pre-remplies. "
              "Repas 50%, assurance non recuperable. PDF telechargeable. "
              "Si methode rapide: toggle par client. Question: 'Ca prend combien de temps a la main?'")

    # Minute 11-14: Audit Module
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "11:00-14:00 --- Module d'audit (si pertinent)")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(1.2), ACCENT_BLUE, [
        ("D'abord, demandez:", True),
        ('"Votre cabinet fait des missions de', False),
        ('verification, d\'examen ou de compilation?"', False),
    ])
    add_box(slide, Inches(6.6), Inches(2.3), Inches(6.1), Inches(1.2), MID_BLUE, [
        ("Si NON:", True),
        ('"C\'est un avantage pour quand vous voudrez', False),
        ('offrir ces services." Passez a la cloture.', False),
    ], title_color=ORANGE)
    audit_steps = [
        ("/engagements", "Liste des engagements en cours"),
        ("/audit/materiality", "CAS 320 --- seuils calcules auto"),
        ("/working_papers", "Dossiers de travail generes du plan comptable"),
        ("Matrice d'assertions", "CAS 500 --- exhaustivite, exactitude, existence, cesure, classement"),
        ("/financial_statements", "Bilan + resultats depuis la balance de verification"),
    ]
    for i, (route, desc) in enumerate(audit_steps):
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.8 + i * 0.6), Inches(12), Inches(0.55))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{route}  "; r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = GOLD; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(17); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "Demandez d'abord s'ils font de l'assurance. Si oui: montrez engagements, materiality, "
              "working papers, assertions, financial statements. Si non: passez vite et dites que c'est disponible.")

    # Minute 14-15: Close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "14:00-15:00 --- Cloture")
    add_accent_line(slide)
    add_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(2.5), MID_BLUE, [
        ('"Voila ce que OtoCPA fait en 15 minutes de demo."', True),
        ('"Imaginez ce que ca fait sur un mois entier de production."', True),
        ("", False),
        ("PAUSE --- laissez le message s'installer", True),
        ("", False),
        ('"Pour [NomDuCabinet] avec [X] clients: [Y]h economisees/mois.', False),
        ('A 150$/h = [Z]$/mois. OtoCPA [forfait] coute [prix]$/mois."', False),
        ("", False),
        ("PAUSE --- laissez les chiffres parler", True),
    ], title_color=WHITE, body_color=LIGHT_BLUE, title_size=Pt(19), body_size=Pt(18))
    add_box(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.8), GREEN, [
        ('"On fait un pilote de 30 jours. Gratuit. 5 clients cette semaine.', True),
        ('Aucune obligation. Aucun paiement. Si ca ne convient pas,', True),
        ('on desinstalle et on reste amis. Est-ce qu\'on commence?"', True),
    ], title_size=Pt(20))
    add_notes(slide, "Resume l'impact. Pause. Chiffres personnalises. Pause. "
              "Proposez le pilote 30 jours. Arretez de parler. ATTENDEZ leur reponse. "
              "Si oui: planifiez l'installation. Si 'on va y penser': proposez un rappel [date]. Si non: laissez la fiche.")

    # Recovery Scenarios
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG2)
    add_logo(slide)
    add_title(slide, "Scenarios de recuperation")
    add_accent_line(slide)
    recoveries = [
        ("Dashboard ne charge pas", '"Un instant, le serveur redemarre."\nLancez autofix.py --- "C\'est le diagnostic auto en action."', ORANGE),
        ("Erreur d'extraction", '"Parfait --- c\'est exactement ce que le systeme devrait faire.\nConfiance basse = revision. L\'equipe corrige, le systeme apprend."', ACCENT_BLUE),
        ("Donnees demo maigres", '"Environnement de demo avec donnees synthetiques.\nVotre vrai environnement: vrais montants et fournisseurs."', LAYER_BLUE),
        ("Question QBO", '"Integration QBO directe via API. Pas de vrai compte en demo,\nmais je peux montrer une ecriture construite."', GREEN),
        ("Question inconnue", '"Excellente question. Je verifie et je vous reviens\npar courriel aujourd\'hui." NE BLUFFEZ JAMAIS.', GOLD),
    ]
    for i, (scenario, response, color) in enumerate(recoveries):
        y = Inches(2.2 + i * 1.0)
        txBox = slide.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(0.9))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{scenario}:  "; r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = response.replace("\n", " | "); r2.font.size = Pt(15); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "5 scenarios de recuperation. Ne paniquez jamais. Chaque 'probleme' peut devenir une demo de fonctionnalite. "
              "Regle d'or: ne bluffez jamais. 'Je verifie et je vous reviens' est toujours valide.")

    # Post-Demo Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_logo(slide)
    add_title(slide, "Checklist post-demo")
    add_accent_line(slide)
    checklist = [
        "Envoyer le courriel de suivi dans les 2 heures (EMAIL_TEMPLATES.md)",
        "Inclure les chiffres specifiques de la demo (clients, heures, ROI)",
        "Si OUI au pilote: planifier l'installation dans les 48 heures",
        "Si 'on va y penser': rappel dans le calendrier a la date convenue",
        "Si NON: ajouter a la liste de contact trimestriel",
        "Logger la demo dans le CRM avec notes sur les douleurs specifiques",
        'Desactiver demo_mode: "demo_mode": false',
    ]
    for i, item in enumerate(checklist):
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(2.3 + i * 0.65), Inches(11.5), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "  "; r1.font.size = Pt(18); r1.font.color.rgb = GREEN; r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(19); r2.font.color.rgb = WHITE; r2.font.name = "Segoe UI"
    add_notes(slide, "Checklist obligatoire apres chaque demo. Le suivi dans les 2 heures est le plus important.")

    out = os.path.join(OUT_DIR, "DEMO_SCRIPT.pptx")
    prs.save(out)
    print(f"  DEMO_SCRIPT.pptx --- {len(prs.slides)} slides")


# ============================================================
# MAIN
# ============================================================
if __name__ == "__main__":
    print("Building all sales presentations...")
    build_sales_script()
    build_objections()
    build_one_pager()
    build_email_templates()
    build_demo_script()
    print("Done! All 5 presentations created.")
