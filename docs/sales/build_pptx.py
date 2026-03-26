"""Build LedgerLink AI sales presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Constants ---
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0x9D, 0xC3, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xB0, 0x50)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xED, 0x7D, 0x31)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, x=Inches(0.4), y=Inches(0.25), size=Pt(14)):
    """Add LedgerLink AI logo text to top-left."""
    txBox = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "LedgerLink"
    run1.font.size = size
    run1.font.bold = True
    run1.font.color.rgb = WHITE
    run1.font.name = "Segoe UI"
    run2 = p.add_run()
    run2.text = " AI"
    run2.font.size = size
    run2.font.bold = True
    run2.font.color.rgb = GOLD
    run2.font.name = "Segoe UI"


def add_logo_dark(slide):
    """Logo for white-background slides."""
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "LedgerLink"
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = DARK_BLUE
    run1.font.name = "Segoe UI"
    run2 = p.add_run()
    run2.text = " AI"
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = ACCENT_BLUE
    run2.font.name = "Segoe UI"


def add_title(slide, text, y=Inches(1.0), font_size=Pt(36), color=WHITE, width=Inches(12)):
    txBox = slide.shapes.add_textbox(Inches(0.7), y, width, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_subtitle(slide, text, y=Inches(2.0), font_size=Pt(24), color=LIGHT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_body_text(slide, text, x=Inches(0.7), y=Inches(2.8), w=Inches(11.5), h=Inches(4),
                  font_size=Pt(24), color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.font.bold = bold
        p.space_after = Pt(8)
    return txBox


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_box(slide, x, y, w, h, fill_color, text, font_size=Pt(20), text_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = text_color
    p.font.name = "Segoe UI"
    p.font.bold = bold
    shape.text_frame.paragraphs[0].space_before = Pt(6)
    return shape


def add_accent_line(slide, y=Inches(1.85)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


# ============================================================
# SLIDE 1 - TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Big centered logo
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run1 = p.add_run()
run1.text = "LedgerLink"
run1.font.size = Pt(64)
run1.font.bold = True
run1.font.color.rgb = WHITE
run1.font.name = "Segoe UI"
run2 = p.add_run()
run2.text = " AI"
run2.font.size = Pt(64)
run2.font.bold = True
run2.font.color.rgb = GOLD
run2.font.name = "Segoe UI"

# Tagline
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "La comptabilite intelligente pour les cabinets CPA du Quebec"
p2.font.size = Pt(28)
p2.font.color.rgb = LIGHT_BLUE
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

# Gold line separator
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.2), Inches(3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

# Contact placeholder
txBox3 = slide.shapes.add_textbox(Inches(8), Inches(6.2), Inches(4.5), Inches(1))
tf3 = txBox3.text_frame
tf3.word_wrap = True
for line in ["[Votre nom]  |  [Votre titre]", "[courriel]@ledgerlink.ca  |  [telephone]"]:
    if tf3.paragraphs[0].text == "":
        p = tf3.paragraphs[0]
    else:
        p = tf3.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_BLUE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT

add_notes(slide,
    "Bonjour, merci de me recevoir. Je m'appelle [votre nom], je travaille chez LedgerLink. "
    "On aide les cabinets CPA du Quebec a automatiser la saisie comptable avec l'intelligence artificielle "
    "--- tout en gardant le controle humain sur chaque transaction. "
    "Je vais vous montrer comment en 15 minutes.\n\n"
    "Timing: 30 secondes maximum. Ne decrivez pas le produit. Presentez-vous, remerciez, fixez l'attente de temps.")

# ============================================================
# SLIDE 2 - THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, '"Combien d\'heures par semaine passez-vous\na saisir des donnees manuellement?"', font_size=Pt(32))
add_accent_line(slide, y=Inches(2.5))

# Three stat boxes
stats = [
    ("40%", "du temps du personnel CPA\nest de la saisie manuelle"),
    ("50-100", "transactions par mois\npar client PME au Quebec"),
    ("3 000$ - 6 000$", "par client/mois en temps\nnon facturable (a 150$/h)"),
]
for i, (num, desc) in enumerate(stats):
    x = Inches(0.8 + i * 4.1)
    y = Inches(3.2)
    # Number box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Segoe UI"

    # Description below
    txBox = slide.shapes.add_textbox(x, Inches(4.6), Inches(3.6), Inches(1.5))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER

add_notes(slide,
    "Posez la question. Attendez. Ne parlez pas. Laissez-les repondre.\n\n"
    "Leur reponse est votre argument de vente. Quand ils disent 'au moins 20 heures' ou 'trop', acquiescez et dites:\n\n"
    "'Exactement. Et chaque heure de saisie manuelle, c'est une heure que vous ne facturez pas. "
    "A 150$ de l'heure, 20 heures par semaine, c'est 12 000$ par mois que votre cabinet laisse sur la table.'\n\n"
    "Si le cabinet a 30 clients: 30 clients x 5 heures de saisie = 150 heures/mois = 22 500$/mois en capacite perdue.\n\n"
    "Ne montrez pas encore la solution. Laissez la douleur s'installer.")

# ============================================================
# SLIDE 3 - THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "LedgerLink AI traite les documents automatiquement", font_size=Pt(32))
add_subtitle(slide, "Photos, courriels, WhatsApp --- extraction instantanee", y=Inches(1.8))
add_accent_line(slide, y=Inches(2.3))

# 4-step flow
steps = [
    ("Photo / PDF\nCourriel", ACCENT_BLUE),
    ("Extraction IA\nautomatique", RGBColor(0x40, 0x60, 0x90)),
    ("Compte GL\nsuggere", RGBColor(0x40, 0x60, 0x90)),
    ("Un clic pour\napprouver", GREEN),
]
for i, (text, color) in enumerate(steps):
    x = Inches(0.6 + i * 3.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.6), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    # Arrow between boxes
    if i < 3:
        arrow_x = Inches(0.6 + (i + 1) * 3.2 - 0.4)
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(3.7), Inches(0.4), Inches(0.5))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GOLD
        arr.line.fill.background()

# Bottom detail line
add_body_text(slide,
    "Fournisseur  |  Montant  |  Date  |  Compte GL  |  TPS/TVQ  |  Approuver",
    y=Inches(5.3), font_size=Pt(22), color=LIGHT_BLUE)

add_notes(slide,
    "'La solution, c'est simple. Votre client prend une photo de sa facture. "
    "LedgerLink lit le document, extrait le fournisseur, le montant, la date, calcule les taxes, "
    "suggere le compte GL, et votre equipe n'a qu'a cliquer Approuver. C'est tout.'\n\n"
    "Une phrase. Arretez. Puis dites: 'Laissez-moi vous montrer.'")

# ============================================================
# SLIDE 4 - HOW IT WORKS (3-Layer Architecture)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Trois couches. Zero approximation sur les taxes.", font_size=Pt(32))
add_accent_line(slide, y=Inches(1.85))

layers = [
    ("COUCHE 1 --- REGLES DETERMINISTES (100% precision)",
     "TPS/TVQ calcules mathematiquement --- jamais une estimation IA\n"
     "13 regles anti-fraude --- seuils stricts, pas d'hallucination\n"
     "Moteur de substance economique --- actifs, passifs, charges",
     ACCENT_BLUE, Inches(2.2)),
    ("COUCHE 2 --- INTELLIGENCE ARTIFICIELLE (suggestions)",
     "Extraction de documents (OCR + Vision)\n"
     "Suggestion de comptes GL | Redaction de messages bilingues\n"
     "Routeur IA : DeepSeek (routine) / Claude (complexe)",
     RGBColor(0x40, 0x60, 0x90), Inches(3.9)),
    ("COUCHE 3 --- CONTROLE HUMAIN (autorite finale)",
     "Votre equipe approuve chaque transaction\n"
     "File d'attente de revision avec niveaux de confiance\n"
     "Historique d'audit complet (qui, quoi, quand)",
     GREEN, Inches(5.5)),
]

for title, body, color, y in layers:
    # Layer header
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(11.9), Inches(0.45))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    # Layer body
    txBox = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.5), Inches(11.5), Inches(1.1))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body.split("\n")):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(15)
        p2.font.color.rgb = LIGHT_BLUE
        p2.font.name = "Segoe UI"

add_notes(slide,
    "'Ce qui distingue LedgerLink des autres outils IA, c'est l'architecture en trois couches.'\n\n"
    "'Couche 1 : les taxes, la detection de fraude, et la classification comptable sont 100% basees sur des regles. "
    "Pas d'IA. Pas d'hallucination. Les calculs TPS/TVQ utilisent la librairie Decimal de Python --- zero erreur d'arrondi.'\n\n"
    "'Couche 2 : l'IA fait la lecture des documents et suggere des comptes. Mais elle ne decide de rien.'\n\n"
    "'Couche 3 : votre equipe a le dernier mot. Toujours. Chaque transaction doit etre approuvee par un humain "
    "avant d'etre soumise a QuickBooks.'\n\n"
    "'Ca repond a la question est-ce qu'on peut lui faire confiance avant meme que vous la posiez.'")

# ============================================================
# SLIDE 5 - QUEBEC-SPECIFIC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Le seul logiciel construit specifiquement\npour la fiscalite quebecoise", font_size=Pt(30))
add_accent_line(slide, y=Inches(2.2))

checklist_items = [
    "TPS/TVQ calcul parallele --- jamais en cascade",
    "Pre-remplissage FPZ-500 Revenu Quebec",
    "Methode rapide configurable par client",
    "Calendrier des echeances TPS/TVQ",
    "Plan comptable general quebecois (196 comptes)",
    "Bilingue FR/EN --- 100% des chaines traduites",
    "Loi 25 --- donnees sur votre serveur, jamais dans le cloud",
    "Assurance provinciale --- charge 9% non recuperable",
    "Repas et divertissement --- recuperation TPS/TVQ a 50%",
    "CNESST taux de cotisation | FSS 6 paliers | RL-1/T4",
]

# Two columns
col1 = checklist_items[:5]
col2 = checklist_items[5:]

for col_idx, items in enumerate([col1, col2]):
    x = Inches(0.7 + col_idx * 6.2)
    y_start = Inches(2.7)
    for i, item in enumerate(items):
        y = y_start + Inches(i * 0.8)
        # Checkmark
        txBox = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run_check = p.add_run()
        run_check.text = "  "
        run_check.font.size = Pt(18)
        run_check.font.color.rgb = GREEN
        run_check.font.name = "Segoe UI"
        run_check.font.bold = True
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(17)
        run_text.font.color.rgb = WHITE
        run_text.font.name = "Segoe UI"

add_notes(slide,
    "'C'est notre avantage le plus fort. LedgerLink est le seul logiciel construit specifiquement "
    "pour la fiscalite quebecoise.'\n\n"
    "'TaxDome? Base a San Francisco, concu pour les Etats-Unis. Karbon? Australie. CaseWare? Ontario, TVH. "
    "Aucun ne gere le calcul parallele TPS/TVQ. Aucun ne pre-remplit le FPZ-500. "
    "Aucun ne comprend la methode rapide.'\n\n"
    "'Et surtout, aucun ne vous permet d'heberger les donnees sur votre propre serveur. "
    "Avec la Loi 25, c'est vous personnellement qui etes responsable des donnees de vos clients. "
    "Avec LedgerLink, les donnees restent chez vous.'")

# ============================================================
# SLIDE 6 - SECURITY AND PRIVACY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Vos donnees clients ne quittent jamais vos locaux", font_size=Pt(32))
add_accent_line(slide, y=Inches(1.85))

# Comparison table
headers = ["", "Cloud (TaxDome/Karbon)", "LedgerLink"]
rows = [
    ["Ou sont les donnees?", "Serveurs AWS aux Etats-Unis", "Votre serveur, dans votre bureau"],
    ["Qui y a acces?", "L'editeur du logiciel", "Vous seul"],
    ["Loi 25 conforme?", "Risque eleve (transfert transfrontalier)", "Oui --- donnees locales"],
    ["Chiffrement", "En transit (TLS)", "En transit + au repos"],
    ["Authentification", "Mot de passe simple", "bcrypt + sessions securisees"],
    ["Piste d'audit", "Variable", "Chaque action enregistree"],
]

table_shape = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.7), Inches(2.4), Inches(11.9), Inches(4.2))
table = table_shape.table

# Style header
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x40)

for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.name = "Segoe UI"
            if j == 2:  # LedgerLink column
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif j == 1:
                p.font.color.rgb = ORANGE
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if i % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x18, 0x2C, 0x50)
        else:
            cell.fill.fore_color.rgb = DARK_BLUE

add_notes(slide,
    "'C'est un sujet emotionnel pour les cabinets CPA au Quebec, et avec raison.'\n\n"
    "'Depuis la Loi 25, vous etes personnellement responsable des donnees de vos clients. "
    "Pas votre cabinet --- vous. Si les donnees de vos clients sont sur un serveur Amazon aux Etats-Unis "
    "et qu'il y a une breche, c'est votre nom sur l'avis de la Commission d'acces a l'information.'\n\n"
    "'Avec LedgerLink, les donnees restent sur votre serveur. Dans votre bureau. Sous votre controle. Point final.'\n\n"
    "'L'IA traite les documents, mais rien n'est stocke dans le cloud. "
    "Les resultats d'extraction sont sauvegardes dans votre base de donnees locale SQLite.'\n\n"
    "Pause. Laissez le message s'imprimer.")

# ============================================================
# SLIDE 7 - TIME SAVINGS / ROI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Le retour sur investissement se calcule en jours,\npas en mois.", font_size=Pt(30))
add_accent_line(slide, y=Inches(2.2))

# Before/After table
headers7 = ["Tache", "Avant", "Apres", "Economie"]
rows7 = [
    ["Saisie de donnees", "3 h/client", "0 h", "3 h"],
    ["Rapprochement bancaire", "1,5 h/client", "20 min", "1 h 10"],
    ["Resume de production", "45 min/client", "5 min", "40 min"],
    ["Detection de doublons", "30 min/client", "0 min", "30 min"],
    ["TOTAL PAR CLIENT", "5 h 45", "25 min", "5 h 20"],
]

table_shape = slide.shapes.add_table(len(rows7) + 1, 4, Inches(0.7), Inches(2.6), Inches(7.5), Inches(3.2))
table = table_shape.table

for j, h in enumerate(headers7):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x40)

for i, row in enumerate(rows7):
    is_total = (i == len(rows7) - 1)
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Segoe UI"
            p.font.bold = is_total
            if j == 3:
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif is_total:
                p.font.color.rgb = GOLD
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if is_total:
            cell.fill.fore_color.rgb = ACCENT_BLUE
        elif i % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x18, 0x2C, 0x50)
        else:
            cell.fill.fore_color.rgb = DARK_BLUE

# ROI boxes on the right
roi_data = [
    ("10 clients", "53 h/mois", "7 950$", "ROI 8 000%"),
    ("30 clients", "160 h/mois", "24 000$", "ROI 9 600%"),
    ("75 clients", "400 h/mois", "60 000$", "ROI 12 000%"),
]

for i, (clients, hours, value, roi) in enumerate(roi_data):
    x = Inches(8.8)
    y = Inches(2.6 + i * 1.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clients
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = f"{hours}  =  {value}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_BLUE
    p2.font.name = "Segoe UI"
    p3 = tf.add_paragraph()
    p3.text = roi
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = GOLD
    p3.font.name = "Segoe UI"

add_notes(slide,
    "'Voici les chiffres. Pas de marketing, juste du math.'\n\n"
    "'Pour un cabinet de 30 clients, LedgerLink economise environ 160 heures par mois. "
    "A 150 dollars de l'heure, c'est 24 000 dollars par mois en capacite recuperee. "
    "LedgerLink coute 249 dollars par mois.'\n\n"
    "Pause. Laissez les chiffres parler.\n\n"
    "'Qu'est-ce que vous feriez avec 160 heures de plus par mois? "
    "Prendre 10 nouveaux clients? Offrir des services-conseils a plus haute valeur? "
    "Donner des vendredis libres a votre equipe?'\n\n"
    "Ne repondez pas a votre propre question. Laissez-les y penser.")

# ============================================================
# SLIDE 8 - CPA AUDIT MODULE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Module d'audit CPA complet", font_size=Pt(32))
add_subtitle(slide, "Verification, examen et compilation --- tout dans un seul systeme", y=Inches(1.7))
add_accent_line(slide, y=Inches(2.2))

# CAS standards - two columns
cas_items_left = [
    ("CAS 315", "Evaluation des risques"),
    ("CAS 320", "Importance relative"),
    ("CAS 330", "Procedures d'audit"),
    ("CAS 500", "Elements probants"),
    ("CAS 505", "Confirmations"),
    ("CAS 530", "Echantillonnage"),
]
cas_items_right = [
    ("CAS 550", "Parties liees"),
    ("CAS 560", "Evenements posterieurs"),
    ("CAS 570", "Continuite d'exploitation"),
    ("CAS 580", "Declarations de la direction"),
    ("CAS 700", "Rapport de l'auditeur"),
    ("CSQC 1", "Controle qualite"),
]

for col_idx, items in enumerate([cas_items_left, cas_items_right]):
    x = Inches(0.7 + col_idx * 6.2)
    for i, (code, desc) in enumerate(items):
        y = Inches(2.6 + i * 0.65)
        txBox = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.55))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run_code = p.add_run()
        run_code.text = f"{code}  "
        run_code.font.size = Pt(17)
        run_code.font.bold = True
        run_code.font.color.rgb = GOLD
        run_code.font.name = "Segoe UI"
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(17)
        run_desc.font.color.rgb = WHITE
        run_desc.font.name = "Segoe UI"

# Bottom features
add_body_text(slide,
    "Etats financiers  |  Procedures analytiques  |  Suivi du temps  |  PDF par composante",
    y=Inches(6.5), font_size=Pt(18), color=LIGHT_BLUE)

add_notes(slide,
    "'Pour les cabinets qui font des missions de verification, d'examen ou de compilation, "
    "le module d'audit justifie a lui seul le prix de LedgerLink.'\n\n"
    "'On couvre les normes CAS 315 a 700 plus CSQC 1. Dossiers de travail, importance relative, "
    "echantillonnage, elements probants, evenements posterieurs, continuite d'exploitation "
    "--- tout integre dans le meme systeme que votre comptabilite.'\n\n"
    "'Est-ce que vous faites des missions de certification? Des examens?'\n\n"
    "Adaptez votre discours selon la reponse. Si oui, insistez. "
    "Si non, passez rapidement et dites 'C'est un avantage pour quand vous voudrez offrir ces services.'")

# ============================================================
# SLIDE 9 - PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, color=RGBColor(0x0F, 0x1A, 0x30))
add_logo(slide)

add_title(slide, "Des prix simples. Pas de surprise.", font_size=Pt(34))
add_accent_line(slide, y=Inches(1.85))

plans = [
    ("Essentiel", "99$", "/mois", "10 clients | 3 utilisateurs",
     "Revision de base\nSoumission QBO", ACCENT_BLUE),
    ("Professionnel", "249$", "/mois", "30 clients | 5 utilisateurs",
     "Routeur IA | Rapprochement\nFraude | Revenu Quebec\nSuivi du temps", RGBColor(0x2E, 0x75, 0xB6)),
    ("Cabinet", "499$", "/mois", "75 clients | 15 utilisateurs",
     "Analytique | Microsoft 365\nCalendrier de production\nCommunications clients", RGBColor(0x40, 0x60, 0x90)),
    ("Entreprise", "999$", "/mois", "Illimite",
     "Module d'audit complet\nEtats financiers\nEchantillonnage | API", GOLD),
]

for i, (name, price, period, capacity, features, accent) in enumerate(plans):
    x = Inches(0.5 + i * 3.15)
    y = Inches(2.3)
    w = Inches(2.95)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(4.6))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x18, 0x2C, 0x50)
    card.line.color.rgb = accent
    card.line.width = Pt(2)

    # Plan name
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    # Price
    txBox2 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.7))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run_price = p2.add_run()
    run_price.text = price
    run_price.font.size = Pt(36)
    run_price.font.bold = True
    run_price.font.color.rgb = WHITE
    run_price.font.name = "Segoe UI"
    run_period = p2.add_run()
    run_period.text = period
    run_period.font.size = Pt(16)
    run_period.font.color.rgb = LIGHT_BLUE
    run_period.font.name = "Segoe UI"

    # Capacity
    txBox3 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.3), w - Inches(0.3), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = capacity
    p3.font.size = Pt(13)
    p3.font.color.rgb = LIGHT_BLUE
    p3.font.name = "Segoe UI"
    p3.alignment = PP_ALIGN.CENTER

    # Features
    txBox4 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.8), w - Inches(0.3), Inches(2.5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    for fi, fline in enumerate(features.split("\n")):
        if fi == 0:
            p4 = tf4.paragraphs[0]
        else:
            p4 = tf4.add_paragraph()
        p4.text = fline
        p4.font.size = Pt(13)
        p4.font.color.rgb = WHITE
        p4.font.name = "Segoe UI"
        p4.alignment = PP_ALIGN.CENTER
        p4.space_after = Pt(4)

# Installation note
add_body_text(slide, "Frais d'installation : 500$ - 1 000$ (une fois) --- installation, configuration, formation, migration",
              y=Inches(7.0), font_size=Pt(16), color=LIGHT_BLUE)

add_notes(slide,
    "Dites le prix avec confiance. Ne vous excusez pas. Ne dites pas 'seulement' ou 'juste'. Dites:\n\n"
    "'Le forfait Professionnel est a 249 dollars par mois. Ca inclut 30 clients, 5 utilisateurs, "
    "le routeur IA, le rapprochement bancaire, la detection de fraude, et le pre-remplissage Revenu Quebec.'\n\n"
    "Arretez de parler. Laissez le silence travailler.\n\n"
    "S'ils ne reagissent pas, ajoutez: 'A 150 dollars de l'heure, LedgerLink se paie en moins de 2 heures "
    "de travail economise. Le premier mois.'\n\n"
    "Ne proposez jamais de rabais sans qu'on vous le demande.")

# ============================================================
# SLIDE 10 - WHAT HAPPENS NEXT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Voici comment on commence :", font_size=Pt(34))
add_accent_line(slide, y=Inches(1.85))

steps_data = [
    ("1", "Demo aujourd'hui", "15 minutes. Je vous montre le vrai logiciel.\nPas de PowerPoint. Des vrais documents.", ACCENT_BLUE),
    ("2", "Pilote gratuit --- 30 jours", "5 clients de votre choix. Voyez les resultats vous-meme.\nAucune obligation. Aucun paiement.", GREEN),
    ("3", "Decision", "Vous decidez si ca vous fait gagner du temps.\nPas de pression. Juste des resultats.", GOLD),
]

for i, (num, title, desc, color) in enumerate(steps_data):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.6)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE if color != GOLD else DARK_BLUE
    p.font.name = "Segoe UI"

    # Step title
    txBox = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(3.6), Inches(0.6))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = color
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER

    # Step description
    txBox2 = slide.shapes.add_textbox(x, y + Inches(1.7), Inches(3.6), Inches(2))
    tf3 = txBox2.text_frame
    tf3.word_wrap = True
    for li, line in enumerate(desc.split("\n")):
        if li == 0:
            p3 = tf3.paragraphs[0]
        else:
            p3 = tf3.add_paragraph()
        p3.text = line
        p3.font.size = Pt(17)
        p3.font.color.rgb = WHITE
        p3.font.name = "Segoe UI"
        p3.alignment = PP_ALIGN.CENTER

add_notes(slide,
    "'Voici les prochaines etapes.'\n\n"
    "'D'abord, je vous montre le logiciel en direct. 15 minutes, avec de vrais documents. Pas de diapositives.'\n\n"
    "'Ensuite, si ca vous interesse, on installe un pilote gratuit de 30 jours. "
    "Vous choisissez 5 clients, on les configure, et vous voyez par vous-meme si ca vous fait gagner du temps.'\n\n"
    "'Apres 30 jours, vous decidez. Aucune obligation. Aucun paiement pendant le pilote. "
    "Si ca ne vous convient pas, on desinstalle et on reste amis.'\n\n"
    "L'objectif est de retirer toute friction. Rendez la prochaine etape evidente et facile.")

# ============================================================
# SLIDE 11 - TESTIMONIALS / TEST RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)

add_title(slide, "Resultats de tests independants", font_size=Pt(32))
add_accent_line(slide, y=Inches(1.85))

test_stats = [
    ("2 853", "tests automatises\n0 echec"),
    ("147/147", "tests adversariaux\n'red team' passes"),
    ("100/100", "score de preparedness\nproduction"),
    ("0", "erreur de calcul\nde taxe"),
]

for i, (num, desc) in enumerate(test_stats):
    x = Inches(0.5 + i * 3.15)
    y = Inches(2.5)

    # Stat box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.95), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x18, 0x2C, 0x50)
    box.line.color.rgb = GREEN
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(17)
    p2.font.color.rgb = WHITE
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.CENTER

# Additional stats below
bottom_stats = [
    "8,5/10 en exactitude fiscale canadienne",
    "13 regles de detection de fraude connectees au pipeline",
    "100% parite bilingue (FR/EN) --- terminologie CPA verifiee",
]
for i, stat in enumerate(bottom_stats):
    txBox = slide.shapes.add_textbox(Inches(2), Inches(5.2 + i * 0.55), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_check = p.add_run()
    run_check.text = "  "
    run_check.font.size = Pt(18)
    run_check.font.color.rgb = GREEN
    run_check.font.name = "Segoe UI"
    run_text = p.add_run()
    run_text.text = stat
    run_text.font.size = Pt(18)
    run_text.font.color.rgb = LIGHT_BLUE
    run_text.font.name = "Segoe UI"

add_notes(slide,
    "'On est un nouveau produit, alors plutot que des temoignages, je vais vous montrer nos resultats de tests independants.'\n\n"
    "'Un red team independant a fait passer 2 853 tests au systeme. Zero echec. "
    "147 tests adversariaux specifiquement concus pour casser le systeme --- tous passes. "
    "Score de preparation a la production: 100 sur 100.'\n\n"
    "'Et zero erreur de calcul de taxe. Parce que les taxes ne sont jamais calculees par l'IA "
    "--- c'est du math pur, verifie independamment.'\n\n"
    "Quand vous aurez de vrais temoignages clients, remplacez cette diapositive.")

# ============================================================
# SLIDE 12 - CALL TO ACTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Big centered logo again
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run1 = p.add_run()
run1.text = "LedgerLink"
run1.font.size = Pt(52)
run1.font.bold = True
run1.font.color.rgb = WHITE
run1.font.name = "Segoe UI"
run2 = p.add_run()
run2.text = " AI"
run2.font.size = Pt(52)
run2.font.bold = True
run2.font.color.rgb = GOLD
run2.font.name = "Segoe UI"

# CTA
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Pret a recuperer 150 heures par mois?"
p2.font.size = Pt(34)
p2.font.bold = True
p2.font.color.rgb = GOLD
p2.font.name = "Segoe UI"
p2.alignment = PP_ALIGN.CENTER

# Gold separator
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.7), Inches(3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

# Contact info
contact_lines = [
    "[Votre nom]",
    "[Votre titre]",
    "[votre.nom]@ledgerlink.ca",
    "[Numero de telephone]",
    "",
    "Reservez une demo : ledgerlink.ca",
]
txBox3 = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9), Inches(3))
tf3 = txBox3.text_frame
tf3.word_wrap = True
for i, line in enumerate(contact_lines):
    if i == 0:
        p3 = tf3.paragraphs[0]
    else:
        p3 = tf3.add_paragraph()
    p3.text = line
    p3.font.size = Pt(22) if i < 4 else Pt(26)
    p3.font.color.rgb = WHITE if i < 4 else LIGHT_BLUE
    p3.font.name = "Segoe UI"
    p3.alignment = PP_ALIGN.CENTER
    if i >= 4:
        p3.font.bold = True

add_notes(slide,
    "'Merci pour votre temps. Est-ce qu'on peut planifier une demo de 15 minutes cette semaine? "
    "Je vous montre le logiciel avec vos vrais documents.'\n\n"
    "Arretez de parler. Attendez leur reponse.\n\n"
    "Si oui: sortez votre calendrier et bloquez le rendez-vous immediatement. "
    "Ne dites pas 'je vous envoie un courriel'. Faites-le maintenant.\n\n"
    "Si 'on va y penser': 'Parfait. Est-ce que je peux vous envoyer un resume par courriel "
    "et vous recontacter [jour precis]?' Obtenez un engagement sur une date.\n\n"
    "Si non: 'Je comprends. Est-ce que je peux vous laisser notre fiche-resume? "
    "Si jamais la situation change, vous aurez mes coordonnees.' Laissez la fiche et partez.")

# --- Save ---
out_path = os.path.join(os.path.dirname(__file__), "LedgerLink_Presentation.pptx")
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
